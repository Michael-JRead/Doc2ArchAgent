#!/usr/bin/env python3
"""Generate Doc2ArchAgent presentation as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors
PRIMARY_DARK = RGBColor(0x1A, 0x1A, 0x2E)
SECONDARY_DARK = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x0F, 0x34, 0x60)
HIGHLIGHT_RED = RGBColor(0xE9, 0x45, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18,
                color=DARK_TEXT, bold=False, font_name='Calibri',
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_bullets(slide, left, top, width, height, items, size=18,
                color=DARK_TEXT, font_name='Calibri', spacing=8):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(spacing)
        p.level = 0
        # bullet character
        p.text = "\u2022  " + item
    return txBox


def add_slide_num(slide, num):
    add_textbox(slide, 12.3, 7.0, 0.8, 0.4, str(num),
                size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def add_title(slide, text, size=32, color=DARK_TEXT):
    add_textbox(slide, 0.8, 0.4, 11.5, 0.8, text, size=size,
                color=color, bold=True)
    # accent line
    add_rect(slide, 0.8, 1.15, 2.0, 0.05, HIGHLIGHT_RED)


def text_in_rect(slide, left, top, width, height, fill_color,
                 title_text, body_text, title_size=20, body_size=16):
    add_rect(slide, left, top, width, height, fill_color)
    add_textbox(slide, left + 0.2, top + 0.2, width - 0.4, 0.5,
                title_text, size=title_size, color=WHITE, bold=True)
    add_textbox(slide, left + 0.2, top + 0.8, width - 0.4, height - 1.0,
                body_text, size=body_size, color=WHITE)


# ============================================================
# SLIDE 1 - Title
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, PRIMARY_DARK)
add_textbox(s, 1.0, 1.5, 11.3, 1.2, "Doc2ArchAgent",
            size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 2.8, 11.3, 0.7,
            "From Documentation to Architecture \u2014 Automatically",
            size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 3.6, 11.3, 0.6,
            "AI-Powered Architecture Extraction, Validation & Visualization",
            size=18, color=HIGHLIGHT_RED, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 5.5, 11.3, 0.5,
            "[Presenter Name]  |  [Date]",
            size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 - The Problem
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "The Problem We All Face")
add_bullets(s, 1.0, 1.6, 10.0, 4.5, [
    "Architecture documentation is scattered across dozens of documents",
    "PDFs, Word docs, Confluence pages, diagrams on whiteboards",
    "Nobody has the complete picture of the system",
    "When someone leaves, critical knowledge walks out the door",
    "Manual architecture reviews take weeks, not hours",
], size=18)
add_textbox(s, 8.5, 6.3, 4.0, 0.5, "Sound familiar?",
            size=22, color=HIGHLIGHT_RED, italic=True, align=PP_ALIGN.RIGHT)
add_slide_num(s, 2)

# ============================================================
# SLIDE 3 - The Cost
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "What This Costs Us")
text_in_rect(s, 0.8, 2.0, 3.6, 3.5, HIGHLIGHT_RED,
             "TIME",
             "Weeks spent manually reviewing\nvendor docs and creating\narchitecture diagrams",
             title_size=26)
text_in_rect(s, 4.85, 2.0, 3.6, 3.5, ACCENT_BLUE,
             "RISK",
             "Security gaps go unnoticed\nbecause nobody can see\nthe full picture",
             title_size=26)
text_in_rect(s, 8.9, 2.0, 3.6, 3.5, SECONDARY_DARK,
             "KNOWLEDGE",
             "Architecture knowledge lives\nin people's heads,\nnot in systems",
             title_size=26)
add_slide_num(s, 3)

# ============================================================
# SLIDE 4 - Our Solution
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Doc2ArchAgent: The Solution")
add_textbox(s, 1.5, 2.2, 10.3, 0.6,
            "Feed in your documents.", size=24, color=DARK_TEXT,
            align=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 2.9, 10.3, 0.6,
            "Get back a complete, validated architecture model",
            size=24, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_textbox(s, 1.5, 3.6, 10.3, 0.6,
            "with diagrams and security analysis.",
            size=24, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_rect(s, 1.5, 5.0, 10.3, 1.0, ACCENT_BLUE)
add_textbox(s, 1.7, 5.15, 9.9, 0.7,
            "15 specialized AI agents working together",
            size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_slide_num(s, 4)

# ============================================================
# SLIDE 5 - How It Works
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "How It Works \u2014 4 Simple Steps")
steps = [
    ("1", "COLLECT", "Drop in your vendor docs (PDFs, Word, images, diagrams)"),
    ("2", "EXTRACT", "AI reads documents and extracts architecture entities"),
    ("3", "VALIDATE", "Automated checks ensure nothing was hallucinated"),
    ("4", "VISUALIZE", "Get diagrams, security reports, and full documentation"),
]
for i, (num, label, desc) in enumerate(steps):
    y = 1.8 + i * 1.2
    add_rect(s, 1.0, y, 0.7, 0.7, ACCENT_BLUE)
    add_textbox(s, 1.0, y + 0.1, 0.7, 0.5, num,
                size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, 2.0, y + 0.05, 2.0, 0.5, label,
                size=22, color=ACCENT_BLUE, bold=True)
    add_textbox(s, 4.0, y + 0.1, 8.0, 0.5, desc,
                size=18, color=DARK_TEXT)
add_slide_num(s, 5)

# ============================================================
# SLIDE 6 - Zero Hallucination
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "The Zero-Hallucination Promise")
add_bullets(s, 1.0, 1.6, 10.5, 3.0, [
    "Every piece of architecture data has a traceable source citation",
    "If the document doesn't say it, it doesn't exist in the model",
    'No guessing. No assuming. No "typical patterns."',
], size=20, spacing=14)
add_rect(s, 0.8, 5.2, 11.7, 1.0, HIGHLIGHT_RED)
add_textbox(s, 1.0, 5.35, 11.3, 0.7,
            "Confidence scores on every entity \u2014 anything below 95% is flagged for human review",
            size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_slide_num(s, 6)

# ============================================================
# SLIDE 7 - What You Get
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "What You Get")
outputs = [
    "Architecture Diagrams \u2014 5 formats (Mermaid, PlantUML, Draw.io, D2, Structurizr)",
    "Security Analysis \u2014 STRIDE threat modeling with CWE mapping",
    "Compliance Mapping \u2014 PCI-DSS, SOC2, GDPR, HIPAA frameworks",
    "Design Documents \u2014 High-Level Design Docs for Confluence or Markdown",
    "Pattern Library \u2014 Reusable templates for future deployments",
]
add_bullets(s, 1.0, 1.6, 11.0, 4.5, outputs, size=20, spacing=14)
add_slide_num(s, 7)

# ============================================================
# SLIDE 8 - Pattern Library
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Build Once, Reuse Everywhere")
add_bullets(s, 1.0, 1.6, 11.0, 3.5, [
    "Network Patterns: Define your topology once (DMZ, App Tier, Data Tier)",
    "Product Patterns: Model a product once (IBM MQ, API Gateway, etc.)",
    "Pop and Swap: Compose any combination into a new deployment in minutes",
], size=20, spacing=14)
add_rect(s, 1.5, 5.3, 10.3, 0.9, ACCENT_BLUE)
add_textbox(s, 1.7, 5.4, 9.9, 0.7,
            'Think of it like LEGO blocks for architecture',
            size=22, color=WHITE, bold=True, italic=True,
            align=PP_ALIGN.CENTER)
add_slide_num(s, 8)

# ============================================================
# SLIDE 9 - Security Built In
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Security Is Not an Afterthought")
add_bullets(s, 1.0, 1.6, 11.0, 3.0, [
    "Automatic threat modeling using STRIDE methodology",
    "Findings mapped to real-world attack patterns (MITRE ATT&CK)",
    "Compliance frameworks checked automatically",
    "Trust boundary analysis across network zones",
], size=20, spacing=12)
add_rect(s, 0.8, 5.3, 11.7, 0.9, HIGHLIGHT_RED)
add_textbox(s, 1.0, 5.4, 11.3, 0.7,
            "10 security findings detected automatically in our demo deployment",
            size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_slide_num(s, 9)

# ============================================================
# SLIDE 10 - The Agents
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "15 Specialized Agents, One Mission")
# 2x2 grid
boxes = [
    (0.8, 2.0, 5.5, 2.2, HIGHLIGHT_RED, "COLLECT & EXTRACT",
     "doc-collector, doc-extractor"),
    (6.9, 2.0, 5.5, 2.2, ACCENT_BLUE, "MODEL & VALIDATE",
     "architect, validator,\npattern-manager, deployer"),
    (0.8, 4.5, 5.5, 2.2, SECONDARY_DARK, "VISUALIZE",
     "diagram-generator + 5 format agents\n(mermaid, plantuml, drawio, d2, structurizr)"),
    (6.9, 4.5, 5.5, 2.2, PRIMARY_DARK, "ANALYZE & DOCUMENT",
     "security-reviewer, doc-writer"),
]
for left, top, w, h, clr, title, body in boxes:
    text_in_rect(s, left, top, w, h, clr, title, body,
                 title_size=20, body_size=16)
add_slide_num(s, 10)

# ============================================================
# SLIDE 11 - Impact
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Real World Impact")
impacts = [
    ("Architecture Reviews:", "Weeks \u2192 Hours"),
    ("Security Analysis:", "Manual Checklists \u2192 Automated STRIDE"),
    ("Documentation:", "Outdated Word Docs \u2192 Living, Validated Models"),
    ("Knowledge Transfer:", "Tribal Knowledge \u2192 Structured, Searchable YAML"),
    ("Compliance:", "Periodic Audits \u2192 Continuous Validation"),
]
for i, (label, value) in enumerate(impacts):
    y = 1.8 + i * 1.0
    add_textbox(s, 1.0, y, 4.5, 0.5, label,
                size=20, color=ACCENT_BLUE, bold=True)
    add_textbox(s, 5.5, y, 7.0, 0.5, value,
                size=20, color=DARK_TEXT)
add_slide_num(s, 11)

# ============================================================
# SLIDE 12 - Security Split (Simple)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Latest Enhancement: Clean Separation")
add_bullets(s, 1.0, 1.6, 11.0, 3.5, [
    'Separated "what the architecture looks like" from "how secure it is"',
    "Diagram pipeline gets clean, focused data \u2014 just structure",
    "Security pipeline gets rich threat modeling data \u2014 CIA, encryption, compliance",
    "7 output files per deployment instead of 3",
], size=20, spacing=12)
add_rect(s, 1.5, 5.3, 10.3, 0.8, ACCENT_BLUE)
add_textbox(s, 1.7, 5.4, 9.9, 0.6,
            "Result: Faster diagrams, deeper security analysis, cleaner data",
            size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_slide_num(s, 12)

# ============================================================
# SLIDE 13 - Section Divider
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, PRIMARY_DARK)
add_textbox(s, 1.0, 2.5, 11.3, 1.2, "Under the Hood",
            size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 3.8, 11.3, 0.7,
            "Technical Deep Dive & Live Demo",
            size=24, color=HIGHLIGHT_RED, align=PP_ALIGN.CENTER)
add_slide_num(s, 13)

# ============================================================
# SLIDE 14 - Architecture Technical
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "C4 Model + Security Layers")
# C4 stack
c4 = [("Context", ACCENT_BLUE), ("Container", SECONDARY_DARK),
      ("Component", SECONDARY_DARK), ("Deployment", PRIMARY_DARK)]
for i, (label, clr) in enumerate(c4):
    y = 2.0 + i * 1.2
    add_rect(s, 1.0, y, 4.5, 0.9, clr)
    add_textbox(s, 1.2, y + 0.15, 4.1, 0.6, label,
                size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
# Security layers
sec = [("Network Zones", ACCENT_BLUE), ("Trust Boundaries", HIGHLIGHT_RED),
       ("Security Analysis", HIGHLIGHT_RED)]
for i, (label, clr) in enumerate(sec):
    y = 2.3 + i * 1.4
    add_rect(s, 7.0, y, 5.0, 1.0, clr)
    add_textbox(s, 7.2, y + 0.2, 4.6, 0.6, label,
                size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 6.5, 5.0, 0.4, "Architecture Layers",
            size=14, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 7.0, 6.5, 5.0, 0.4, "Security Layers",
            size=14, color=HIGHLIGHT_RED, bold=True, align=PP_ALIGN.CENTER)
add_slide_num(s, 14)

# ============================================================
# SLIDE 15 - YAML Model
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Structured Data, Not Documents")
# Left column
add_textbox(s, 1.0, 1.8, 5.0, 0.5, "Base Files (Diagrams)",
            size=20, color=ACCENT_BLUE, bold=True)
for i, f in enumerate(["system.yaml", "networks.yaml", "deployment.yaml"]):
    add_textbox(s, 1.5, 2.5 + i * 0.6, 4.0, 0.5, f,
                size=18, color=DARK_TEXT, font_name='Consolas')
# Right column
add_textbox(s, 7.0, 1.8, 5.5, 0.5, "Security Overlays (Threats)",
            size=20, color=HIGHLIGHT_RED, bold=True)
for i, f in enumerate(["system-security.yaml", "networks-security.yaml",
                         "deployment-security.yaml"]):
    add_textbox(s, 7.5, 2.5 + i * 0.6, 5.0, 0.5, f,
                size=18, color=DARK_TEXT, font_name='Consolas')
# Divider
add_rect(s, 6.3, 1.8, 0.04, 2.5, LIGHT_GRAY)
# Bottom
add_textbox(s, 1.0, 5.0, 11.3, 0.5,
            "+ provenance.yaml (source citations)  |  All validated against JSON schemas",
            size=15, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
add_slide_num(s, 15)

# ============================================================
# SLIDE 16 - Security Split Technical
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Security YAML Split \u2014 Pipeline Separation")
# Left box
add_rect(s, 0.8, 2.0, 5.5, 3.5, ACCENT_BLUE)
add_textbox(s, 1.0, 2.1, 5.1, 0.5, "DIAGRAM PIPELINE",
            size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bullets(s, 1.2, 2.7, 4.8, 2.5, [
    "system.yaml", "networks.yaml", "deployment.yaml",
    "Clean, focused, structural data",
], size=16, color=WHITE, spacing=10)
# Right box
add_rect(s, 7.0, 2.0, 5.5, 3.5, HIGHLIGHT_RED)
add_textbox(s, 7.2, 2.1, 5.1, 0.5, "THREAT PIPELINE",
            size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_bullets(s, 7.4, 2.7, 4.8, 2.5, [
    "All 7 files merged by entity ID",
    "merge_security_overlays()",
    "Zero changes to rule engine",
    "Additive merge before construction",
], size=16, color=WHITE, spacing=10)
# Bottom stats
add_textbox(s, 0.8, 6.0, 11.7, 0.5,
            "60-68% of fields in each file were security-only  |  366 tests pass",
            size=15, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
add_slide_num(s, 16)

# ============================================================
# SLIDE 17 - Compose Pipeline
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Pop and Swap Composition")
# Flow boxes
add_rect(s, 1.0, 2.2, 3.5, 1.5, SECONDARY_DARK)
add_textbox(s, 1.2, 2.3, 3.1, 0.6, "manifest.yaml",
            size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 1.2, 2.9, 3.1, 0.6, "1 network + N products",
            size=16, color=WHITE, align=PP_ALIGN.CENTER)
# Arrow
add_rect(s, 4.8, 2.7, 2.5, 0.06, ACCENT_BLUE)
add_textbox(s, 4.8, 2.2, 2.5, 0.5, "compose.py",
            size=16, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
# Output
add_rect(s, 7.5, 2.2, 4.5, 1.5, ACCENT_BLUE)
add_textbox(s, 7.7, 2.3, 4.1, 0.6, "7 Output Files",
            size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 7.7, 2.9, 4.1, 0.6, "+ diagrams/ directory",
            size=16, color=WHITE, align=PP_ALIGN.CENTER)
# Bullets
add_bullets(s, 1.0, 4.3, 11.0, 2.5, [
    "ID prefixes prevent conflicts between patterns",
    "Cross-product relationships supported",
    "Security fields auto-separated into overlay files",
], size=18, spacing=10)
add_slide_num(s, 17)

# ============================================================
# SLIDE 18 - Threat Rules
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Automated Threat Detection")
add_bullets(s, 1.0, 1.6, 11.0, 2.0, [
    "Deterministic rule engine \u2014 same input, same findings, every time",
    "STRIDE per element with CWE mapping",
    "Severity scoring with accepted risk tracking",
], size=18, spacing=10)
# Code-style box
add_rect(s, 0.8, 4.0, 11.7, 2.5, PRIMARY_DARK)
add_textbox(s, 1.0, 4.1, 11.3, 0.4, "Example Findings:",
            size=14, color=LIGHT_GRAY, bold=True)
findings = [
    "HIGH   | missing-logging-zone         | prod-dmz",
    "MED    | no-input-validation          | console-to-qm",
    "LOW    | no-supply-chain-attestation   | mq-console",
    "LOW    | missing-monitoring-zone       | prod-private-app-tier",
    "INFO   | missing-data-classification   | console-to-qm",
]
for i, f in enumerate(findings):
    add_textbox(s, 1.2, 4.5 + i * 0.38, 10.8, 0.35, f,
                size=14, color=WHITE, font_name='Consolas')
add_slide_num(s, 18)

# ============================================================
# SLIDE 19 - Demo Agenda
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Live Demo")
demo_steps = [
    "Create a deployment manifest",
    "Run compose.py to generate all 7 files",
    "Run validate.py with security overlays",
    "Run threat-rules.py to detect findings",
    "Inspect the composed output files",
    "Walk through a security finding",
]
for i, step in enumerate(demo_steps):
    y = 1.8 + i * 0.8
    add_rect(s, 1.0, y, 0.6, 0.6, ACCENT_BLUE)
    add_textbox(s, 1.0, y + 0.08, 0.6, 0.45, str(i + 1),
                size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(s, 2.0, y + 0.1, 9.0, 0.45, step,
                size=20, color=DARK_TEXT)
add_slide_num(s, 19)

# ============================================================
# SLIDE 20 - Demo Commands
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, PRIMARY_DARK)
add_textbox(s, 0.8, 0.4, 11.5, 0.7, "Demo Commands",
            size=30, color=WHITE, bold=True)
cmds = [
    "python tools/compose.py deployments/mq-prod-us-east/manifest.yaml --validate",
    "python tools/validate.py system.yaml networks.yaml \\\n"
    "  --security system-security.yaml --networks-security networks-security.yaml",
    "python tools/threat-rules.py system.yaml --networks networks.yaml \\\n"
    "  --security system-security.yaml --format table",
    "python -m pytest tests/ -v   # 366 passed, 0 skipped",
]
for i, cmd in enumerate(cmds):
    y = 1.5 + i * 1.4
    add_rect(s, 0.8, y, 11.7, 1.1, SECONDARY_DARK)
    add_textbox(s, 1.0, y + 0.15, 11.3, 0.8, cmd,
                size=14, color=WHITE, font_name='Consolas')
add_slide_num(s, 20)

# ============================================================
# SLIDE 21 - Q&A
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, PRIMARY_DARK)
add_textbox(s, 1.0, 2.0, 11.3, 1.2, "Questions?",
            size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 3.8, 11.3, 0.7,
            "github.com/Michael-JRead/Doc2ArchAgent",
            size=22, color=HIGHLIGHT_RED, align=PP_ALIGN.CENTER)
add_textbox(s, 1.0, 5.5, 11.3, 0.5,
            "[Your Name]  |  [your.email@company.com]",
            size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 22 - Appendix
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, LIGHT_BG)
add_title(s, "Appendix: Test Coverage")
add_textbox(s, 1.0, 1.6, 11.0, 0.7, "366 tests  |  0 skipped",
            size=30, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_bullets(s, 1.0, 2.6, 11.0, 4.5, [
    "Schema validity \u2014 all 12 JSON schemas verified",
    "Example conformance \u2014 all example YAML files validated",
    "Security split \u2014 15+ tests for overlay merge, cross-reference, coverage",
    "Compose pipeline \u2014 7-file output, prefixing, context merging",
    "Threat rules engine \u2014 finding detection, severity, SARIF output",
    "Every schema, every tool, every agent doc \u2014 regression tested",
], size=18, spacing=12)
add_slide_num(s, 22)

# ============================================================
# SAVE
# ============================================================
out_path = os.path.join(os.path.dirname(__file__),
                        "Doc2ArchAgent-Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Size: {os.path.getsize(out_path) / 1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
